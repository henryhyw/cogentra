import csv
import hashlib
import io
import json
import re
import time
import uuid
import zipfile
from collections import Counter
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from docx import Document
from pptx import Presentation
from pypdf import PdfReader
from sqlalchemy.orm import Session

from oralv.models import (
    Artifact,
    ArtifactKind,
    ArtifactStatus,
    AssessmentCase,
    CaseStatus,
    CompetencyState,
    CompetencyStatus,
    DefensePlan,
    DefenseQuestion,
    DefenseSession,
    EvidenceItem,
    EvidenceStatus,
    GenerationStage,
    GenerationTrace,
    ParsedArtifact,
    PlanStatus,
    QuestionKind,
    ResponseSegment,
    ResponseStatus,
    SessionStatus,
    SubmissionRecord,
    TranscriptSegment,
)
from oralv.providers.registry import registry
from oralv.services.audit import log_event
from oralv.storage import StorageClient


@dataclass
class ParsedContent:
    text: str
    normalized: dict[str, Any]
    low_text_confidence: bool


def classify_artifact(filename: str, content_type: str) -> tuple[ArtifactKind, float]:
    stem = filename.lower()
    if "rubric" in stem or "criteria" in stem:
        return ArtifactKind.rubric, 0.95
    if "assignment" in stem or "brief" in stem or "instruction" in stem:
        return ArtifactKind.assignment, 0.92
    if "reference" in stem or "notes" in stem:
        return ArtifactKind.reference, 0.86
    if content_type.startswith("video/") or content_type.startswith("audio/"):
        return ArtifactKind.recording, 0.97
    if stem.endswith(".txt") and "transcript" in stem:
        return ArtifactKind.transcript, 0.93
    return ArtifactKind.submission, 0.82


def infer_student_identity(filename: str) -> tuple[str, str]:
    stem = Path(filename).stem.replace("_", " ").replace("-", " ").strip()
    student_name = re.sub(r"\s+", " ", stem).title() or "Student Submission"
    local = re.sub(r"[^a-z0-9]+", ".", student_name.lower()).strip(".")
    return student_name, f"{local or 'student'}@demo.student"


def _extract_pdf(data: bytes) -> str:
    reader = PdfReader(io.BytesIO(data))
    parts = [page.extract_text() or "" for page in reader.pages]
    return "\n\n".join(parts)


def _extract_docx(data: bytes) -> str:
    document = Document(io.BytesIO(data))
    return "\n".join(paragraph.text for paragraph in document.paragraphs if paragraph.text.strip())


def _extract_pptx(data: bytes) -> str:
    presentation = Presentation(io.BytesIO(data))
    slides: list[str] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        fragments = [f"Slide {slide_index}"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                fragments.append(shape.text)
        slides.append("\n".join(fragments))
    return "\n\n".join(slides)


def _extract_zip(data: bytes) -> str:
    parts: list[str] = []
    with zipfile.ZipFile(io.BytesIO(data)) as archive:
        for entry in archive.infolist():
            if entry.is_dir():
                continue
            if not entry.filename.endswith((".py", ".md", ".txt", ".json", ".ipynb", ".csv")):
                continue
            raw = archive.read(entry.filename)
            parts.append(f"# {entry.filename}\n{raw.decode('utf-8', errors='ignore')}")
    return "\n\n".join(parts)


def extract_text(filename: str, content_type: str, data: bytes) -> str:
    suffix = Path(filename).suffix.lower()
    if suffix in {".txt", ".md", ".py", ".json", ".csv"}:
        return data.decode("utf-8", errors="ignore")
    if suffix == ".pdf":
        return _extract_pdf(data)
    if suffix == ".docx":
        return _extract_docx(data)
    if suffix == ".pptx":
        return _extract_pptx(data)
    if suffix == ".zip":
        return _extract_zip(data)
    if content_type.startswith("text/"):
        return data.decode("utf-8", errors="ignore")
    return ""


def _split_sections(text: str) -> list[dict[str, Any]]:
    blocks = [chunk.strip() for chunk in re.split(r"\n\s*\n", text) if chunk.strip()]
    sections: list[dict[str, Any]] = []
    for index, block in enumerate(blocks[:24], start=1):
        title = block.splitlines()[0][:80]
        sections.append({"id": f"section-{index}", "title": title, "content": block[:1200]})
    return sections


def _extract_rubric_criteria(text: str) -> list[str]:
    criteria: list[str] = []
    for line in text.splitlines():
        normalized = line.strip(" -\t")
        if not normalized:
            continue
        if ":" in normalized or len(normalized.split()) <= 12:
            criteria.append(normalized[:140])
    deduped = list(dict.fromkeys(criteria))
    return deduped[:8]


def _extract_claims(text: str) -> list[str]:
    sentences = re.split(r"(?<=[.!?])\s+", text)
    claims = [
        sentence.strip()
        for sentence in sentences
        if len(sentence.split()) >= 8 and any(word in sentence.lower() for word in ["because", "therefore", "designed", "used", "proposed", "evaluated", "implemented"])
    ]
    if not claims:
        claims = [sentence.strip() for sentence in sentences if len(sentence.split()) >= 10]
    return claims[:6]


def normalize_artifact(kind: ArtifactKind, family: str, text: str) -> ParsedContent:
    sections = _split_sections(text)
    low_conf = len(text.strip()) < 120
    payload: dict[str, Any] = {
        "family": family,
        "summary": text[:500],
        "sections": sections,
    }
    if kind == ArtifactKind.rubric:
        payload["criteria"] = _extract_rubric_criteria(text)
    if kind == ArtifactKind.submission:
        payload["claims"] = _extract_claims(text)
        payload["keywords"] = [word for word, _ in Counter(re.findall(r"[A-Za-z]{5,}", text.lower())).most_common(12)]
    if kind == ArtifactKind.assignment:
        payload["objectives"] = _extract_rubric_criteria(text)[:5]
    return ParsedContent(text=text, normalized=payload, low_text_confidence=low_conf)


def _generation_trace(
    db: Session,
    *,
    organization_id,
    case_id,
    submission_record_id,
    defense_session_id,
    stage: GenerationStage,
    prompt_name: str,
    payload: dict[str, Any],
    output_entity_type: str,
    output_entity_id: str,
) -> GenerationTrace:
    llm = registry.llm()
    started = time.perf_counter()
    result = llm.generate(prompt_name, "v1", payload)
    elapsed = int((time.perf_counter() - started) * 1000)
    trace = GenerationTrace(
        organization_id=organization_id,
        case_id=case_id,
        submission_record_id=submission_record_id,
        defense_session_id=defense_session_id,
        stage=stage,
        provider=llm.name,
        model_name=llm.model_name,
        prompt_name=prompt_name,
        prompt_version="v1",
        input_hash=hashlib.sha256(json.dumps(payload, sort_keys=True).encode()).hexdigest(),
        seed=result.metadata["seed"],
        latency_ms=elapsed,
        output_entity_type=output_entity_type,
        output_entity_id=output_entity_id,
        metadata_json=result.metadata,
    )
    db.add(trace)
    db.flush()
    return trace


def parse_artifact(db: Session, artifact: Artifact, case: AssessmentCase) -> ParsedArtifact:
    storage = StorageClient()
    raw = storage.download_bytes(artifact.storage_key)
    extracted = extract_text(artifact.filename, artifact.content_type, raw)
    normalized = normalize_artifact(artifact.kind, case.assignment_family.value, extracted)
    parsed = (
        db.query(ParsedArtifact)
        .filter(ParsedArtifact.artifact_id == artifact.id)
        .first()
    )
    if not parsed:
        parsed = ParsedArtifact(
            organization_id=artifact.organization_id,
            case_id=artifact.case_id,
            artifact_id=artifact.id,
            submission_record_id=artifact.submission_record_id,
        )
        db.add(parsed)
    parsed.normalized_text = normalized.text
    parsed.normalized_json = normalized.normalized
    parsed.section_count = len(normalized.normalized.get("sections", []))
    parsed.word_count = len(normalized.text.split())
    parsed.parsing_confidence = 0.42 if normalized.low_text_confidence else 0.88
    parsed.low_text_confidence = normalized.low_text_confidence
    parsed.embedding = registry.embedding().embed(normalized.text[:2000])
    artifact.status = ArtifactStatus.ready
    db.flush()
    _generation_trace(
        db,
        organization_id=artifact.organization_id,
        case_id=artifact.case_id,
        submission_record_id=artifact.submission_record_id,
        defense_session_id=None,
        stage=GenerationStage.normalization,
        prompt_name="normalize-artifact",
        payload={"artifact_id": str(artifact.id), "kind": artifact.kind.value, "text": normalized.text[:2000]},
        output_entity_type="parsed_artifact",
        output_entity_id=str(parsed.id),
    )
    return parsed


def ensure_submission_for_artifact(db: Session, artifact: Artifact) -> SubmissionRecord | None:
    if artifact.kind != ArtifactKind.submission:
        return None
    if artifact.submission_record_id:
        return db.query(SubmissionRecord).filter(SubmissionRecord.id == artifact.submission_record_id).first()
    student_name, student_email = infer_student_identity(artifact.filename)
    queue_position = (
        db.query(SubmissionRecord).filter(SubmissionRecord.case_id == artifact.case_id).count() + 1
    )
    submission = SubmissionRecord(
        organization_id=artifact.organization_id,
        case_id=artifact.case_id,
        student_name=student_name,
        student_email=student_email,
        external_identifier=Path(artifact.filename).stem.lower().replace(" ", "-"),
        status="processing",
        queue_position=queue_position,
        artifact_bundle_json={"source_artifact_ids": [str(artifact.id)]},
    )
    db.add(submission)
    db.flush()
    artifact.submission_record_id = submission.id
    return submission


def build_plan_for_submission(db: Session, case: AssessmentCase, submission: SubmissionRecord) -> DefensePlan:
    assignment = (
        db.query(ParsedArtifact)
        .join(Artifact, ParsedArtifact.artifact_id == Artifact.id)
        .filter(Artifact.case_id == case.id, Artifact.kind == ArtifactKind.assignment)
        .first()
    )
    rubric = (
        db.query(ParsedArtifact)
        .join(Artifact, ParsedArtifact.artifact_id == Artifact.id)
        .filter(Artifact.case_id == case.id, Artifact.kind == ArtifactKind.rubric)
        .first()
    )
    submission_parsed = (
        db.query(ParsedArtifact)
        .filter(ParsedArtifact.submission_record_id == submission.id)
        .order_by(ParsedArtifact.created_at.asc())
        .first()
    )
    claims = submission_parsed.normalized_json.get("claims", []) if submission_parsed else []
    criteria = rubric.normalized_json.get("criteria", []) if rubric else []
    objectives = assignment.normalized_json.get("objectives", []) if assignment else []
    focus_points: list[dict[str, Any]] = []
    base_targets = criteria or objectives or ["understanding", "method", "judgment"]
    for index, target in enumerate(base_targets[:5], start=1):
        linked_claim = claims[(index - 1) % len(claims)] if claims else submission_parsed.normalized_text[:180]
        focus_points.append(
            {
                "id": f"focus-{index}",
                "target": target,
                "claim": linked_claim,
                "risk": "Needs oral elaboration" if len(linked_claim.split()) < 18 else "Requires evidence under questioning",
            }
        )

    plan = (
        db.query(DefensePlan)
        .filter(DefensePlan.submission_record_id == submission.id)
        .order_by(DefensePlan.version_number.desc())
        .first()
    )
    next_version = 1 if not plan else plan.version_number + 1
    plan = DefensePlan(
        organization_id=case.organization_id,
        case_id=case.id,
        submission_record_id=submission.id,
        status=PlanStatus.ready,
        version_number=next_version,
        focus_graph_json={"nodes": focus_points, "family": case.assignment_family.value},
        plan_json={
            "claims": claims,
            "criteria": criteria,
            "objectives": objectives,
            "summary": submission_parsed.normalized_json.get("summary") if submission_parsed else "",
        },
        settings_json={"attempts": 1, "prep_seconds": 45, "response_seconds": 150},
    )
    db.add(plan)
    db.flush()

    for index, focus in enumerate(focus_points, start=1):
        prompt = (
            f"Walk me through how your submission demonstrates {focus['target']}. "
            f"Refer to the specific choice or claim: {focus['claim'][:140]}"
        )
        question = DefenseQuestion(
            defense_plan_id=plan.id,
            sequence=index,
            prompt=prompt,
            rationale=f"Tests whether the student can defend the submission against {focus['target']}.",
            evaluation_target=focus["target"],
            expected_signal=f"Explains reasoning, tradeoffs, and evidence related to {focus['target']}.",
            question_kind=QuestionKind.primary,
            prep_seconds=45,
            response_seconds=150,
            branch_condition_json={"follow_up_when": "confidence_low_or_generic_answer"},
            embedding=registry.embedding().embed(prompt),
        )
        db.add(question)
    db.flush()
    _generation_trace(
        db,
        organization_id=case.organization_id,
        case_id=case.id,
        submission_record_id=submission.id,
        defense_session_id=None,
        stage=GenerationStage.question_generation,
        prompt_name="generate-defense-plan",
        payload={"focus_points": focus_points, "submission_id": str(submission.id)},
        output_entity_type="defense_plan",
        output_entity_id=str(plan.id),
    )
    log_event(
        db,
        organization_id=case.organization_id,
        event_type="plan.generated",
        entity_type="defense_plan",
        entity_id=str(plan.id),
        actor_type="system",
        case_id=case.id,
        submission_record_id=submission.id,
        payload={"version": next_version},
    )
    return plan


def process_case_ingestion(db: Session, case_id: str) -> dict[str, Any]:
    case = db.query(AssessmentCase).filter(AssessmentCase.id == case_id).first()
    if not case:
        return {"status": "missing"}
    case.status = CaseStatus.ingesting
    artifacts = db.query(Artifact).filter(Artifact.case_id == case.id, Artifact.deleted_at.is_(None)).all()
    for artifact in artifacts:
        if artifact.kind == ArtifactKind.submission and artifact.submission_record_id is None:
            ensure_submission_for_artifact(db, artifact)
        parse_artifact(db, artifact, case)
    submissions = (
        db.query(SubmissionRecord)
        .filter(SubmissionRecord.case_id == case.id, SubmissionRecord.deleted_at.is_(None))
        .order_by(SubmissionRecord.queue_position.asc())
        .all()
    )
    for submission in submissions:
        build_plan_for_submission(db, case, submission)
        submission.status = "ready_for_plan_review"
    case.status = CaseStatus.ready_for_review
    log_event(
        db,
        organization_id=case.organization_id,
        event_type="case.processed",
        entity_type="assessment_case",
        entity_id=str(case.id),
        actor_type="system",
        case_id=case.id,
    )
    db.commit()
    return {"status": "processed", "submission_count": len(submissions)}


def publish_session_for_plan(db: Session, plan: DefensePlan) -> DefenseSession:
    existing = (
        db.query(DefenseSession)
        .filter(DefenseSession.defense_plan_id == plan.id)
        .order_by(DefenseSession.created_at.desc())
        .first()
    )
    if existing:
        return existing
    session = DefenseSession(
        organization_id=plan.organization_id,
        case_id=plan.case_id,
        submission_record_id=plan.submission_record_id,
        defense_plan_id=plan.id,
        public_token=f"sess_{uuid.uuid4().hex[:20]}",
        access_token_hash=hashlib.sha256(uuid.uuid4().hex.encode()).hexdigest(),
        status=SessionStatus.invited,
        settings_snapshot_json=plan.settings_json,
    )
    db.add(session)
    plan.status = PlanStatus.published
    case = db.query(AssessmentCase).filter(AssessmentCase.id == plan.case_id).first()
    if case:
        case.status = CaseStatus.published
        case.published_at = case.published_at or __import__("datetime").datetime.now(__import__("datetime").UTC)
    db.flush()
    log_event(
        db,
        organization_id=plan.organization_id,
        event_type="session.published",
        entity_type="defense_session",
        entity_id=str(session.id),
        actor_type="system",
        case_id=plan.case_id,
        submission_record_id=plan.submission_record_id,
        defense_session_id=session.id,
    )
    db.commit()
    return session


def _segment_sentences(text: str) -> list[str]:
    parts = [segment.strip() for segment in re.split(r"(?<=[.!?])\s+", text) if segment.strip()]
    return parts[:8] or [text]


def synthesize_evidence(db: Session, defense_session_id: str) -> dict[str, Any]:
    session = db.query(DefenseSession).filter(DefenseSession.id == defense_session_id).first()
    if not session:
        return {"status": "missing"}
    plan = db.query(DefensePlan).filter(DefensePlan.id == session.defense_plan_id).first()
    questions = (
        db.query(DefenseQuestion)
        .filter(DefenseQuestion.defense_plan_id == session.defense_plan_id)
        .order_by(DefenseQuestion.sequence.asc())
        .all()
    )
    responses = (
        db.query(ResponseSegment)
        .filter(ResponseSegment.defense_session_id == session.id)
        .order_by(ResponseSegment.sequence.asc())
        .all()
    )
    existing_evidence = db.query(EvidenceItem).filter(EvidenceItem.defense_session_id == session.id).all()
    for item in existing_evidence:
        db.delete(item)
    existing_competencies = (
        db.query(CompetencyStatus).filter(CompetencyStatus.submission_record_id == session.submission_record_id).all()
    )
    for item in existing_competencies:
        db.delete(item)
    db.flush()

    competency_rollup: dict[str, list[EvidenceItem]] = {}
    for question, response in zip(questions, responses, strict=False):
        transcript = response.transcript_text or ""
        transcript_lower = transcript.lower()
        score = 0
        if len(transcript.split()) > 30:
            score += 1
        if any(token in transcript_lower for token in ["because", "tradeoff", "evidence", "example", "decision"]):
            score += 1
        if question.evaluation_target.lower() in transcript_lower:
            score += 1
        status = EvidenceStatus.verified if score >= 2 else EvidenceStatus.unresolved if score == 1 else EvidenceStatus.inconsistent
        evidence = EvidenceItem(
            organization_id=session.organization_id,
            case_id=session.case_id,
            submission_record_id=session.submission_record_id,
            defense_session_id=session.id,
            defense_question_id=question.id,
            competency_key=re.sub(r"[^a-z0-9]+", "-", question.evaluation_target.lower()).strip("-"),
            title=question.evaluation_target.title(),
            status=status,
            evidence_snippet=transcript[:280],
            rationale="Response was analyzed against question target, specificity, and supporting detail.",
            confidence=0.42 + score * 0.18,
            related_claim=plan.plan_json.get("claims", [""])[0] if plan.plan_json.get("claims") else None,
            source_segment_ids_json=[],
        )
        db.add(evidence)
        db.flush()
        competency_rollup.setdefault(evidence.competency_key, []).append(evidence)

    for key, items in competency_rollup.items():
        label = items[0].title
        if any(item.status == EvidenceStatus.inconsistent for item in items):
            state = CompetencyState.inconsistent
        elif any(item.status == EvidenceStatus.unresolved for item in items):
            state = CompetencyState.unresolved
        else:
            state = CompetencyState.verified
        competency = CompetencyStatus(
            organization_id=session.organization_id,
            case_id=session.case_id,
            submission_record_id=session.submission_record_id,
            competency_key=key,
            label=label,
            status=state,
            score=sum(float(item.confidence) for item in items) / len(items),
            summary=f"{label} is currently assessed as {state.value.replace('_', ' ')} based on the recorded defense.",
            supporting_evidence_ids_json=[str(item.id) for item in items],
        )
        db.add(competency)

    session.status = SessionStatus.completed
    session.completed_at = __import__("datetime").datetime.now(__import__("datetime").UTC)
    submission = db.query(SubmissionRecord).filter(SubmissionRecord.id == session.submission_record_id).first()
    if submission:
        submission.status = "ready_for_decision"
    case = db.query(AssessmentCase).filter(AssessmentCase.id == session.case_id).first()
    if case:
        case.status = CaseStatus.reviewing
    _generation_trace(
        db,
        organization_id=session.organization_id,
        case_id=session.case_id,
        submission_record_id=session.submission_record_id,
        defense_session_id=session.id,
        stage=GenerationStage.evidence_alignment,
        prompt_name="synthesize-evidence",
        payload={"session_id": str(session.id), "response_count": len(responses)},
        output_entity_type="defense_session",
        output_entity_id=str(session.id),
    )
    log_event(
        db,
        organization_id=session.organization_id,
        event_type="evidence.synthesized",
        entity_type="defense_session",
        entity_id=str(session.id),
        actor_type="system",
        case_id=session.case_id,
        submission_record_id=session.submission_record_id,
        defense_session_id=session.id,
        payload={"evidence_count": sum(len(items) for items in competency_rollup.values())},
    )
    db.commit()
    return {"status": "completed", "competencies": len(competency_rollup)}


def to_csv(rows: list[dict[str, Any]]) -> str:
    if not rows:
        return ""
    buffer = io.StringIO()
    writer = csv.DictWriter(buffer, fieldnames=list(rows[0].keys()))
    writer.writeheader()
    writer.writerows(rows)
    return buffer.getvalue()
